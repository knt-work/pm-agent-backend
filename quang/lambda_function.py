# lambda_function.py
# -*- coding: utf-8 -*-

import json
import os
import io

# boto3 optional for local runs
try:
    import boto3  # type: ignore
    s3_client = boto3.client("s3")
except Exception:
    boto3 = None
    s3_client = None

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
import openpyxl  # read embedded Excel from charts

# -----------------------------------------------------------------------------
# Utilities
# -----------------------------------------------------------------------------

def compact(obj):
    """
    Recursively remove keys with None / False / empty list/dict.
    Preserves True and non-empty values.
    """
    if isinstance(obj, dict):
        out = {}
        for k, v in obj.items():
            cv = compact(v)
            if cv in (None, False, [], {}):
                continue
            out[k] = cv
        return out
    if isinstance(obj, list):
        out = [compact(v) for v in obj]
        out = [v for v in out if v not in (None, False, [], {})]
        return out
    return obj

# -----------------------------------------------------------------------------
# Heuristics for roles + serializers
# -----------------------------------------------------------------------------

def _shape_role_for_text(shape):
    """Heuristic role detection for text shapes."""
    try:
        if getattr(shape, "is_placeholder", False):
            pht = shape.placeholder_format.type
            if pht in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return "title"
            if pht == PP_PLACEHOLDER.SUBTITLE:
                return "subtitle"
    except Exception:
        pass

    # Font-size heuristic (first run with size)
    try:
        if shape.has_text_frame and shape.text_frame.paragraphs:
            for run in shape.text_frame.paragraphs[0].runs:
                if run.font.size:
                    pt = getattr(run.font.size, "pt", None)
                    if pt is not None:
                        if pt >= 28:
                            return "title"
                        if pt >= 20:
                            return "heading"
    except Exception:
        pass
    return "body"


def _paragraph_dict(paragraph):
    """Serialize a paragraph minimally."""
    d = {
        "level": paragraph.level or 0,
        "text": paragraph.text
    }
    # alignment (only if meaningful)
    try:
        if isinstance(paragraph.alignment, PP_ALIGN) and paragraph.alignment is not None:
            d["alignment"] = str(paragraph.alignment)
    except Exception:
        pass
    return compact(d)


def _table_dict(table):
    """Serialize table grid, keeping only r/c/text and merge flags when true."""
    rows = []
    for r_idx, row in enumerate(table.rows):
        row_cells = []
        for c_idx, cell in enumerate(row.cells):
            text_val = ""
            try:
                if getattr(cell, "text_frame", None):
                    text_val = cell.text_frame.text
                else:
                    text_val = cell.text
            except Exception:
                text_val = ""
            cell_obj = {"r": r_idx, "c": c_idx, "text": text_val}
            # Include merge flags only when True
            if bool(getattr(cell, "is_merge_origin", False)):
                cell_obj["merge_origin"] = True
            if bool(getattr(cell, "is_spanned", False)):
                cell_obj["is_spanned"] = True
            row_cells.append(cell_obj)
        rows.append(row_cells)

    try:
        n_rows = len(table.rows)
    except Exception:
        n_rows = None
    try:
        n_cols = len(table.columns)
    except Exception:
        n_cols = None

    return compact({
        "rows": n_rows,
        "cols": n_cols,
        "cells": rows
    })

# -----------------------------------------------------------------------------
# Extractors (compact)
# -----------------------------------------------------------------------------

def extract_chart_data(chart):
    """Structured chart (compact)."""
    data = {
        "kind": "chart",
        "title": chart.chart_title.text_frame.text if getattr(chart, "has_title", False) else None,
        "chart_type": getattr(getattr(chart, "chart_type", None), "name", str(getattr(chart, "chart_type", ""))),
        "series": [],
        "excel_data": None
    }

    # Series + points
    try:
        for s in chart.series:
            pts = []
            cats = getattr(s, "categories", None)
            vals = getattr(s, "values", None)
            if cats is not None and vals is not None:
                for i in range(min(len(cats), len(vals))):
                    # category label
                    cat_obj = cats[i]
                    cat_label = None
                    try:
                        cat_label = getattr(cat_obj, "label", None)
                    except Exception:
                        pass
                    if not cat_label:
                        try:
                            cat_label = str(cat_obj)
                        except Exception:
                            cat_label = None

                    val = vals[i]
                    try:
                        val = float(val) if val is not None else None
                    except Exception:
                        pass
                    pts.append(compact({"category": cat_label, "value": val}))
            data["series"].append(compact({
                "name": getattr(s, "name", None),
                "points": pts
            }))
    except Exception:
        # fallback
        data["series"] = [compact({"name": getattr(s, "name", None)}) for s in getattr(chart, "series", [])]

    # Embedded Excel (best-effort)
    try:
        chart_part = chart.part
        embedded_excel_blob = chart_part.chart_workbook.xlsx_part.blob
        workbook = openpyxl.load_workbook(io.BytesIO(embedded_excel_blob), data_only=True, read_only=True)
        sheet = workbook.active
        grid = []
        for row in sheet.iter_rows(values_only=True):
            grid.append(list(row))
        data["excel_data"] = grid
    except Exception:
        pass

    return compact(data)


def extract_text_shape(shape):
    """Compact text node: role + paragraphs (level, text)."""
    role = _shape_role_for_text(shape)
    paragraphs = []
    for p in shape.text_frame.paragraphs:
        paragraphs.append(_paragraph_dict(p))
    node = {
        "kind": "text",
        "role": role,
        "paragraphs": paragraphs
    }
    return compact(node)


def extract_table_shape(shape):
    """Compact table node with grid."""
    table_info = _table_dict(shape.table)
    node = {
        "kind": "table",
        "structure": table_info
    }
    return compact(node)


def extract_picture_shape(shape):
    """Very small image node."""
    content_type = None
    try:
        content_type = getattr(shape.image, "content_type", None)
    except Exception:
        pass
    return compact({
        "kind": "image",
        "content_type": content_type
    })

# -----------------------------------------------------------------------------
# Core analyzer (shared by Lambda and local runner)
# -----------------------------------------------------------------------------

def _analyze_presentation_stream(file_stream, file_key="<in-memory>"):
    prs = Presentation(file_stream)
    file_results = {
        "file_name": file_key,
        "slide_count": len(prs.slides),
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": i + 1,
            "elements": []
        }
        elements = []

        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    elements.append(extract_chart_data(shape.chart))
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    elements.append(extract_table_shape(shape))
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    elements.append(extract_picture_shape(shape))
                elif getattr(shape, "has_text_frame", False):
                    elements.append(extract_text_shape(shape))
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # Keep a group with compact children
                    group_children = []
                    for s in shape.shapes:
                        try:
                            if s.shape_type == MSO_SHAPE_TYPE.CHART:
                                group_children.append(extract_chart_data(s.chart))
                            elif s.shape_type == MSO_SHAPE_TYPE.TABLE:
                                group_children.append(extract_table_shape(s))
                            elif s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                group_children.append(extract_picture_shape(s))
                            elif getattr(s, "has_text_frame", False):
                                group_children.append(extract_text_shape(s))
                        except Exception:
                            continue
                    elements.append(compact({"kind": "group", "children": group_children}))
                # else: ignore unknowns altogether to keep output lean
            except Exception:
                continue

        # Keep natural order (no z-order metadata stored)
        slide_data["elements"] = compact(elements)
        file_results["slides"].append(slide_data)

    return compact(file_results)

# -----------------------------------------------------------------------------
# Lambda handler
# -----------------------------------------------------------------------------

def lambda_handler(event, context):
    BUCKET_NAME = os.environ.get('BUCKET_NAME', 'default-bucket')

    try:
        body = json.loads(event.get('body', '{}')) if isinstance(event, dict) else {}
        file_keys = body.get('fileKeys')

        if not file_keys:
            return {'statusCode': 400, 'body': json.dumps({'message': "'fileKeys' (array) is required."})}

        if s3_client is None:
            return {'statusCode': 500, 'body': json.dumps({'message': "S3 client not available in this environment."})}

        full_analysis = {}
        processed_files = []
        failed_files = []

        for file_key in file_keys:
            if not file_key.endswith('.pptx'):
                continue

            try:
                response = s3_client.get_object(Bucket=BUCKET_NAME, Key=file_key)
                file_stream = io.BytesIO(response['Body'].read())
                result = _analyze_presentation_stream(file_stream, file_key)
                full_analysis[file_key] = result
                processed_files.append(file_key)
            except Exception as e:
                print(f"Error processing file {file_key}: {e}")
                failed_files.append(file_key)

        # Compact summary
        summary_text = f"Processed {len(processed_files)} files. "
        for key, value in full_analysis.items():
            total_text = sum(1 for s in value['slides'] for e in s['elements'] if e.get('kind') == 'text')
            total_tables = sum(1 for s in value['slides'] for e in s['elements'] if e.get('kind') == 'table')
            total_charts = sum(1 for s in value['slides'] for e in s['elements'] if e.get('kind') == 'chart')
            summary_text += (
                f"File '{key}' ({value['slide_count']} slides): "
                f"{total_text} text blocks, {total_tables} tables, {total_charts} charts. "
            )

        return {
            'statusCode': 200,
            'headers': {'Content-Type': 'application/json', 'Access-Control-Allow-Origin': '*'},
            'body': json.dumps({
                'message': 'Analysis complete.',
                'summary': summary_text,
                'full_analysis_snippet': str(full_analysis)[:1000] + "...",
                'processed': processed_files,
                'failed': failed_files
            }, ensure_ascii=False)
        }

    except Exception as e:
        print(f"Handler error: {e}")
        return {
            'statusCode': 500,
            'headers': {'Content-Type': 'application/json', 'Access-Control-Allow-Origin': '*'},
            'body': json.dumps({'message': 'Error processing files'})
        }
