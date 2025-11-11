from __future__ import annotations
import os
from datetime import datetime, timedelta
from dateutil import parser as dateparser

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# =========================
# Settings
# =========================
WIDESCREEN_16x9 = True           # keep 16:9 by default
TITLE = "Task Roadmap (Q1’25 → Q1’26)"
T0 = "2025-01-01"
T1 = "2026-03-31"

LANE_COUNT = 15
LANE_NAMES = [f"Task {i+1}" for i in range(LANE_COUNT)]  # Task 1..Task 15 (Task 1 at TOP)

# Make the chevron point slimmer (≈0.15–0.5 is reasonable; smaller = slimmer point)
CHEVRON_HEAD = 0.28

# Bright, readable palette
PALETTE = [
    "#4FC3F7", "#29B6F6", "#81C784", "#66BB6A", "#FFD54F",
    "#FFB74D", "#FF8A65", "#BA68C8", "#9575CD", "#64B5F6",
    "#4DB6AC", "#AED581", "#F06292", "#FFCC80", "#90CAF9"
]

def parse_time(x):
    if isinstance(x, (int, float)):
        return datetime.fromtimestamp(x)
    return dateparser.parse(x)

def to_rgb(hex_color: str) -> RGBColor:
    s = hex_color.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

def script_dir() -> str:
    return os.path.dirname(os.path.abspath(__file__))

def generate_items():
    """Create one chevron bar per lane with staggered dates and names."""
    t0 = parse_time(T0)
    items = []
    for i, lane_name in enumerate(LANE_NAMES):
        start_dt = t0 + timedelta(days=8 + i * 22)       # staggered starts
        duration = 90 + (i * 9 % 80)                     # 90..170 days
        end_dt = start_dt + timedelta(days=duration)

        items.append({
            "id": f"task-{i+1}",
            "kind": "task",
            "label": f"{lane_name} – Phase {chr(65 + (i % 3))}",
            "start": {"x": start_dt.date().isoformat(), "y": i, "z": 0},
            "end":   {"x": end_dt.date().isoformat()},
            "size":  {"height": 0.62},
            "style": {"fill": PALETTE[i % len(PALETTE)], "text": "#1B1B1B"}  # dark text on bright fill
        })
    return items

# SPEC = {
#     "version": "1.0",
#     "units": {
#         "x": "time",
#         "y": "lanes",
#         "z": "layers",
#         "xRange": {"t0": T0, "t1": T1},
#         "yRange": {"lanes": LANE_NAMES}
#     },
#     "legend": {"items": []},
#     "axisPoints": [],  # milestones omitted (clean look)
#     "items": generate_items(),
#     "metadata": {"title": TITLE}
# }

SPEC = {
  "version": "1.0",
  "units": {
    "x": "time",
    "y": "lanes",
    "z": "layers",
    "xRange": { "t0": "2025-01-01", "t1": "2026-12-31" },
    "yRange": {
      "lanes": [
        "Phase 1: Research & Planning",
        "Phase 2: Architecture & System Design",
        "Phase 3: Core Backend Development",
        "Phase 4: Frontend MVP Build",
        "Phase 5: UX/UI Refinement",
        "Phase 6: Integration Development",
        "Phase 7: Analytics & Reporting",
        "Phase 8: System Testing & QA",
        "Phase 9: Public Launch Preparation",
        "Phase 10: Feedback & Iteration"
      ]
    }
  },
  "legend": { "items": [] },
  "axisPoints": [],
  "items": [
    { "id": "task-1", "kind": "task", "label": "Phase 1: Research & Planning", "start": { "x": "2025-01-01", "y": 0, "z": 0 }, "end": { "x": "2025-03-31" }, "size": { "height": 0.62 }, "style": { "fill": "#4FC3F7", "text": "#1B1B1B" } },
    { "id": "task-2", "kind": "task", "label": "Phase 2: Architecture & System Design", "start": { "x": "2025-01-01", "y": 1, "z": 0 }, "end": { "x": "2025-06-30" }, "size": { "height": 0.62 }, "style": { "fill": "#29B6F6", "text": "#1B1B1B" } },
    { "id": "task-3", "kind": "task", "label": "Phase 3: Core Backend Development", "start": { "x": "2025-04-01", "y": 2, "z": 0 }, "end": { "x": "2025-06-30" }, "size": { "height": 0.62 }, "style": { "fill": "#81C784", "text": "#1B1B1B" } },
    { "id": "task-4", "kind": "task", "label": "Phase 4: Frontend MVP Build", "start": { "x": "2025-04-01", "y": 3, "z": 0 }, "end": { "x": "2025-09-30" }, "size": { "height": 0.62 }, "style": { "fill": "#66BB6A", "text": "#1B1B1B" } },
    { "id": "task-5", "kind": "task", "label": "Phase 5: UX/UI Refinement", "start": { "x": "2025-07-01", "y": 4, "z": 0 }, "end": { "x": "2025-09-30" }, "size": { "height": 0.62 }, "style": { "fill": "#FFD54F", "text": "#1B1B1B" } },
    { "id": "task-6", "kind": "task", "label": "Phase 6: Integration Development", "start": { "x": "2025-07-01", "y": 5, "z": 0 }, "end": { "x": "2025-12-31" }, "size": { "height": 0.62 }, "style": { "fill": "#FFB74D", "text": "#1B1B1B" } },
    { "id": "task-7", "kind": "task", "label": "Phase 7: Analytics & Reporting", "start": { "x": "2025-10-01", "y": 6, "z": 0 }, "end": { "x": "2025-12-31" }, "size": { "height": 0.62 }, "style": { "fill": "#FF8A65", "text": "#1B1B1B" } },
    { "id": "task-8", "kind": "task", "label": "Phase 8: System Testing & QA", "start": { "x": "2026-01-01", "y": 7, "z": 0 }, "end": { "x": "2026-03-31" }, "size": { "height": 0.62 }, "style": { "fill": "#BA68C8", "text": "#1B1B1B" } },
    { "id": "task-9", "kind": "task", "label": "Phase 9: Public Launch Preparation", "start": { "x": "2026-04-01", "y": 8, "z": 0 }, "end": { "x": "2026-06-30" }, "size": { "height": 0.62 }, "style": { "fill": "#9575CD", "text": "#1B1B1B" } },
    { "id": "task-10", "kind": "task", "label": "Phase 10: Feedback & Iteration", "start": { "x": "2026-07-01", "y": 9, "z": 0 }, "end": { "x": "2026-12-31" }, "size": { "height": 0.62 }, "style": { "fill": "#64B5F6", "text": "#1B1B1B" } }
  ],
  "metadata": { "title": "Product Roadmap (2025–2026)" }
}


def build_ppt_from_spec(spec: dict, output_filename: str = "diagram_output.pptx") -> str:
    prs = Presentation()
    # Force 16:9 unless you flip the flag
    if WIDESCREEN_16x9:
        prs.slide_width  = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width  = Inches(10.0)
        prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide_width, slide_height = prs.slide_width, prs.slide_height

    # Layout
    margin_left, margin_top, margin_right, margin_bottom = Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0)
    canvas_left = margin_left
    canvas_top  = margin_top + Inches(0.35)  # room for Q labels
    canvas_width  = slide_width - margin_left - margin_right
    canvas_height = slide_height - canvas_top - margin_bottom

    lanes = spec["units"]["yRange"]["lanes"]
    lane_count = len(lanes)
    lane_height = canvas_height / max(lane_count, 1)

    # X mapping
    t0 = parse_time(spec["units"]["xRange"]["t0"])
    t1 = parse_time(spec["units"]["xRange"]["t1"])
    total_days = max((t1 - t0).days, 1)

    def x_pos(dt_str: str):
        dt = parse_time(dt_str)
        d = (dt - t0).days
        return canvas_left + canvas_width * (d / total_days)

    def y_pos(lane_index: int):
        # Task 1 at TOP (no inversion)
        return canvas_top + lane_height * lane_index

    # Title
    if (title := spec.get("metadata", {}).get("title")):
        tb = slide.shapes.add_textbox(margin_left, Inches(0.35) - Inches(0.25), canvas_width, Inches(0.4))
        tf = tb.text_frame
        tf.text = title
        p = tf.paragraphs[0]; p.font.size = Pt(22); p.font.bold = True

    # Alternating lane backgrounds
    for i in range(lane_count):
        y = y_pos(i)
        if i % 2 == 0:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, canvas_left, y, canvas_width, lane_height)
            bg.fill.solid(); bg.fill.fore_color.rgb = to_rgb("#F3F4F6")
            bg.line.fill.background()

    # Quarter grid & labels
    q_boundaries = ["2025-01-01","2025-04-01","2025-07-01","2025-10-01","2026-01-01","2026-04-01"]
    q_labels     = [("Q1","2025-01-01"),("Q2","2025-04-01"),("Q3","2025-07-01"),("Q4","2025-10-01"),("Q5","2026-01-01")]

    for d in q_boundaries:
        x = x_pos(d)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, canvas_top, Inches(0.018), canvas_height)
        line.fill.solid(); line.fill.fore_color.rgb = to_rgb("#E0E0E0")
        line.line.fill.background()

    for q, d in q_labels:
        x = x_pos(d)
        tb = slide.shapes.add_textbox(x + Inches(0.05), canvas_top - Inches(0.30), Inches(0.6), Inches(0.25))
        tf = tb.text_frame; tf.text = q
        p = tf.paragraphs[0]; p.font.size = Pt(12)

    # Left gutter labels (Task 1 at TOP)
    for i, name in enumerate(lanes):
        tb = slide.shapes.add_textbox(margin_left - Inches(0.95), y_pos(i), Inches(0.9), lane_height)
        tf = tb.text_frame
        tf.text = name
        p = tf.paragraphs[0]; p.font.size = Pt(12); p.font.bold = True

    # ---- Chevron helper (single shape) ----
    def draw_chevron(left, top, width, height, fill_hex, text, text_color="#1B1B1B"):
        """
        Draw a single CHEVRON shape (one piece, like your image).
        CHEVRON_HEAD controls the point proportion (0.0..1.0).
        """
        chev = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)

        # adjust the "notch" (how far the point cuts in). Try/except for safety on older Office versions.
        try:
            chev.adjustments[0] = CHEVRON_HEAD  # default ~0.5; smaller = slimmer point
        except Exception:
            pass

        chev.fill.solid(); chev.fill.fore_color.rgb = to_rgb(fill_hex)
        chev.line.fill.background()

        tf = chev.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size = Pt(11); p.font.bold = True
        p.font.color.rgb = to_rgb(text_color)
        tf.word_wrap = True
        tf.margin_left = Inches(0.08); tf.margin_top = Inches(0.02)

        return chev

    # Draw chevron bars
    drawn_items = 0
    for it in spec.get("items", []):
        if it.get("kind") != "task":
            continue

        left = x_pos(it["start"]["x"])
        right = x_pos(it["end"]["x"])
        if right < left:
            left, right = right, left
        width = right - left

        lane_idx = int(it["start"]["y"])
        height = lane_height * float(it["size"].get("height", 0.62))
        top = y_pos(lane_idx) + (lane_height - height) / 2

        draw_chevron(left, top, width, height,
                     it["style"].get("fill", "#90CAF9"),
                     it.get("label", it["id"]),
                     it["style"].get("text", "#1B1B1B"))
        drawn_items += 1

    # Save
    out_path = os.path.join(script_dir(), output_filename)
    prs.save(out_path)
    print(f"✅ Saved PowerPoint: {os.path.abspath(out_path)}")
    print(f"   Drawn chevron bars: {drawn_items}")
    print(f"   Slide size: {prs.slide_width} x {prs.slide_height} EMUs (16:9 default)")
    return out_path

if __name__ == "__main__":
    build_ppt_from_spec(SPEC)
