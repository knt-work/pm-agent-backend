# -*- coding: utf-8 -*-
"""
Spec → PPTX renderer (multi-slide, cover slide, auto-fit text, safe legends, gantt gutter)

Elements:
- text: heading | paragraph | bullets | rich
- table
- chart: gantt (chevrons), pie, bar (clustered/stacked, vertical/horizontal), line (markers/smooth)

Samples at the bottom: SAMPLE_TEN_SLIDES (10 slides), SAMPLE_FLAT_WITH_SLIDE_TAGS
"""
from __future__ import annotations
import re
_Q_RE = re.compile(r"^(\d{4})\s+Q([1-4])$", re.I)
import json

def cleaning_JSON(raw_text: str):
    """
    - Parse JSON string to Python objects
    - Convert string 'True'/'False' to Python booleans True/False
      (JSON booleans True/False are already handled by json.loads)
    - Recursively remove any key 'layout' whose value is None.
    """
    data = json.loads(raw_text)

    def normalize(obj):
        # Convert string booleans like "True" / "False"
        if isinstance(obj, str):
            lower = obj.strip().lower()
            if lower == "True":
                return True
            if lower == "False":
                return False
            return obj  # keep as-is

        # Dict: recurse + drop layout=None
        if isinstance(obj, dict):
            new = {}
            for k, v in obj.items():
                if k == "layout" and v is None:
                    continue
                new[k] = normalize(v)
            return new

        # List: recurse
        if isinstance(obj, list):
            return [normalize(v) for v in obj]

        # Other types unchanged (including actual bools from JSON: True/False)
        return obj

    return normalize(data)

import os
from datetime import datetime
try:
    from dateutil import parser as dateparser  # optional but recommended
except Exception:
    dateparser = None

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import json

# charts
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# =========================
# Global Settings
# =========================
WIDESCREEN_16x9 = True
DEFAULT_FONT_NAME = "Calibri"
DEFAULT_TEXT_COLOR = "#1B1B1B"
DEFAULT_BG = "#FFFFFF"

# =========================
# Utils
# =========================

def parse_time(x):
    # already a timestamp
    if isinstance(x, (int, float)):
        return datetime.fromtimestamp(x)

    # already a datetime
    if isinstance(x, datetime):
        return x

    # None or undefined-ish → just return "now" as a safe fallback
    if x is None:
        return datetime.now()

    # handle strings
    if isinstance(x, str):
        s = x.strip()

        # special case: "2025 Q1" → 2025-01-01, "2025 Q2" → 2025-04-01, etc.
        m = _Q_RE.match(s)
        if m:
            year = int(m.group(1))
            q = int(m.group(2))
            month = (q - 1) * 3 + 1
            return datetime(year, month, 1)

        # try dateutil
        if dateparser:
            try:
                return dateparser.parse(s)
            except Exception:
                pass

        # try ISO8601
        try:
            return datetime.fromisoformat(s)
        except Exception:
            # last-ditch fallback: don't crash, just return "now"
            return datetime.now()

    # very defensive fallback for weird types
    s = str(x)
    if dateparser:
        try:
            return dateparser.parse(s)
        except Exception:
            return datetime.now()
    try:
        return datetime.fromisoformat(s)
    except Exception:
        return datetime.now()



def to_rgb(hex_color: str) -> RGBColor:
    s = (hex_color or "").lstrip("#")
    if len(s) == 3:
        s = "".join([c*2 for c in s])
    if len(s) != 6:
        s = "1B1B1B"
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

def script_dir() -> str:
    return os.path.dirname(os.path.abspath(__file__))

def inches_rect(position):
    x = Inches(position.get("x", 1.0))
    y = Inches(position.get("y", 1.0))
    w = Inches(position.get("w", 10.0))
    h = Inches(position.get("h", 1.0))
    return x, y, w, h

def apply_text_style(paragraph, style: dict | None):
    if not style:
        style = {}
    font = getattr(paragraph, "font", None)
    if font is None:
        font = paragraph
    f = style.get("font", {})
    font.name = f.get("name", DEFAULT_FONT_NAME)
    if "size" in f and f["size"]:
        font.size = Pt(f["size"])
    if "bold" in f:
        font.bold = bool(f["bold"])
    if "italic" in f:
        font.italic = bool(f["italic"])
    if "underline" in f:
        font.underline = bool(f["underline"])
    color = style.get("text", DEFAULT_TEXT_COLOR)
    font.color.rgb = to_rgb(color)
    align = style.get("align")
    if align and hasattr(paragraph, "alignment"):
        paragraph.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT
        }.get(align, PP_ALIGN.LEFT)

# =========================
# Shared chart helpers
# =========================
def _style_chart_legend(chart, position="right"):
    """Place legend and reserve layout space so it won't overlap the plot."""
    chart.has_legend = True
    pos = {
        "right": XL_LEGEND_POSITION.RIGHT,
        "left": XL_LEGEND_POSITION.LEFT,
        "top": XL_LEGEND_POSITION.TOP,
        "bottom": XL_LEGEND_POSITION.BOTTOM,
        "corner": XL_LEGEND_POSITION.CORNER,
    }.get(position, XL_LEGEND_POSITION.RIGHT)
    chart.legend.position = pos
    chart.legend.include_in_layout = True  # <-- prevent overlap

# =========================
# Renderers: TEXT
# =========================
def render_textbox(slide, element):
    """
    Supports: heading | paragraph | bullets | rich
    Auto-fit enabled; optional margins: style.margins.{left,right,top,bottom} in inches

    If element["position"] is not provided, sensible defaults are chosen per variant
    that work nicely on a 16:9 slide.
    """
    variant = element.get("variant", "paragraph")
    style = element.get("style", {}) or {}

    # -------- Default positions (if not explicitly provided) --------
    if "position" in element:
        position = element["position"]
    else:
        # These assume a 13.333" wide slide and leave ~1" margins.
        if variant == "heading":
            # Big title near the top
            position = {"x": 1.0, "y": 0.8, "w": 11.3, "h": 1.2}
        elif variant == "bullets":
            # Bullets under a typical heading
            position = {"x": 1.0, "y": 1.9, "w": 11.3, "h": 4.0}
        else:  # paragraph / rich
            position = {"x": 1.0, "y": 1.4, "w": 11.3, "h": 3.0}

    x, y, w, h = inches_rect(position)

    # -------- Create textbox & frame --------
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    # Margins
    margins = style.get("margins", {}) or {}
    try:
        if "left" in margins:
            tf.margin_left = Inches(margins["left"])
        if "right" in margins:
            tf.margin_right = Inches(margins["right"])
        if "top" in margins:
            tf.margin_top = Inches(margins["top"])
        if "bottom" in margins:
            tf.margin_bottom = Inches(margins["bottom"])
    except Exception:
        pass

    # -------- Variant-specific rendering --------
    if variant in ("heading", "paragraph"):
        p = tf.paragraphs[0]
        p.text = element.get("text", "")

        # For heading, enforce a bold, larger font, but still allow
        # caller to override size/name if they provided them.
        if variant == "heading":
            font_cfg = style.get("font", {})
            font_cfg = {
                "name": font_cfg.get("name", DEFAULT_FONT_NAME),
                "size": font_cfg.get("size", 28),
                "bold": True,
            }
            style = {
                "font": font_cfg,
                "align": style.get("align", "left"),
                "text": style.get("text", DEFAULT_TEXT_COLOR),
                "margins": style.get("margins", {}),
            }

        apply_text_style(p, style)

    elif variant == "bullets":
        items = element.get("items", []) or []
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item.get("text", "")
            p.level = int(item.get("level", 0))
            apply_text_style(p, style)

            # Optional per-bullet font size override
            font_cfg = style.get("font", {}) or {}
            if "size" in font_cfg:
                p.font.size = Pt(font_cfg["size"])
            else:
                p.font.size = Pt(16)

    elif variant == "rich":
        runs = element.get("runs", []) or []
        p = tf.paragraphs[0]
        p.text = ""
        for r in runs:
            run = p.add_run()
            run.text = r.get("text", "")
            rfont = r.get("font", {}) or {}
            run.font.name = rfont.get("name", DEFAULT_FONT_NAME)
            if "size" in rfont:
                run.font.size = Pt(rfont["size"])
            if "bold" in r:
                run.font.bold = bool(r["bold"])
            if "italic" in r:
                run.font.italic = bool(r["italic"])
            if "underline" in r:
                run.font.underline = bool(r["underline"])
            run.font.color.rgb = to_rgb(
                r.get("color", style.get("text", DEFAULT_TEXT_COLOR))
            )

        apply_text_style(p, style)


# =========================
# Renderer: TABLE
# =========================
def render_table(slide, element):
    headers = element.get("headers", [])
    rows = element.get("rows", [])
    ncols = max(len(headers), max((len(r) for r in rows), default=0))
    nrows = 1 + len(rows)

    x, y, w, h = inches_rect(element.get("position", {"x": 1, "y": 1, "w": 10, "h": 2}))
    shape = slide.shapes.add_table(nrows, ncols, x, y, w, h)
    table = shape.table

    col_widths = element.get("column_widths")
    if col_widths:
        for i, width in enumerate(col_widths[:ncols]):
            table.columns[i].width = Inches(width)

    # header
    for c in range(ncols):
        cell = table.cell(0, c)
        cell.text = headers[c] if c < len(headers) else ""
        p = cell.text_frame.paragraphs[0]
        apply_text_style(p, {"font": {"size": 12, "bold": True}, "align": element.get("style", {}).get("align", "left")})

    # body
    for r, row in enumerate(rows, start=1):
        for c in range(ncols):
            cell = table.cell(r, c)
            cell.text = row[c] if c < len(row) else ""
            p = cell.text_frame.paragraphs[0]
            apply_text_style(p, {"font": {"size": 12}, "align": element.get("style", {}).get("align", "left")})

# =========================
# Charts: Pie / Bar / Line
# =========================
def render_chart_pie(slide, element):
    legend_pos = element.get("legend", "right")
    legend_pad = float(element.get("legendPadInches", 0.7))
    x, y, w, h = _padded_chart_frame(
        element.get("position", {"x": 6, "y": 2, "w": 5, "h": 4}),
        legend_pos,
        legend_pad,
    )

    data = element.get("data", [])
    chart_data = ChartData()
    chart_data.categories = [d["label"] for d in data]
    chart_data.add_series(element.get("title", ""), [d["value"] for d in data])

    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, w, h, chart_data).chart

    # Title
    if element.get("title"):
        chart.has_title = True
        chart.chart_title.text_frame.text = element["title"]

    # Legend (non-overlapping)
    _style_chart_legend(chart, legend_pos)

    # Slice colors
    series = chart.series[0]
    for i, d in enumerate(data):
        col = d.get("color")
        if col:
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = to_rgb(col)

    # Optional labels (%)
    try:
        plot = chart.plots[0]
        if element.get("showLabels"):
            plot.has_data_labels = True
            if element.get("labels") == "percent":
                plot.data_labels.number_format = "0%"
                plot.data_labels.show_percentage = True
    except Exception:
        pass


def render_chart_bar(slide, element):
    cats = element.get("x", {}).get("categories", [])
    series = element.get("series", [])
    opts = element.get("options", {})
    stacked = bool(opts.get("stacked", False))
    orientation = opts.get("orientation", "vertical")

    if orientation == "horizontal":
        chart_type = (XL_CHART_TYPE.BAR_STACKED if stacked
                      else XL_CHART_TYPE.BAR_CLUSTERED)
    else:
        chart_type = (XL_CHART_TYPE.COLUMN_STACKED if stacked
                      else XL_CHART_TYPE.COLUMN_CLUSTERED)

    legend_pos = element.get("legend", "right")
    legend_pad = float(element.get("legendPadInches", 0.7))
    x, y, w, h = _padded_chart_frame(
        element.get("position", {"x": 1, "y": 2, "w": 11, "h": 4}),
        legend_pos,
        legend_pad,
    )

    chart_data = CategoryChartData()
    chart_data.categories = cats
    for s in series:
        chart_data.add_series(s.get("name", ""), s.get("data", []))

    chart = slide.shapes.add_chart(chart_type, x, y, w, h, chart_data).chart

    if element.get("title"):
        chart.has_title = True
        chart.chart_title.text_frame.text = element["title"]

    _style_chart_legend(chart, legend_pos)

    # Series colors
    for i, s in enumerate(series):
        col = s.get("color")
        if col:
            chart.series[i].format.fill.solid()
            chart.series[i].format.fill.fore_color.rgb = to_rgb(col)


def render_chart_line(slide, element):
    cats = element.get("x", {}).get("categories", [])
    series = element.get("series", [])
    opts = element.get("options", {})
    markers = bool(opts.get("markers", True))
    smooth = bool(opts.get("smooth", False))

    chart_type = XL_CHART_TYPE.LINE_MARKERS if markers else XL_CHART_TYPE.LINE

    legend_pos = element.get("legend", "right")
    legend_pad = float(element.get("legendPadInches", 0.7))
    x, y, w, h = _padded_chart_frame(
        element.get("position", {"x": 1, "y": 2, "w": 11, "h": 4}),
        legend_pos,
        legend_pad,
    )

    chart_data = CategoryChartData()
    chart_data.categories = cats
    for s in series:
        chart_data.add_series(s.get("name", ""), s.get("data", []))

    chart = slide.shapes.add_chart(chart_type, x, y, w, h, chart_data).chart

    if element.get("title"):
        chart.has_title = True
        chart.chart_title.text_frame.text = element["title"]

    _style_chart_legend(chart, legend_pos)

    # Per-series styling
    for i, s in enumerate(series):
        ser = chart.series[i]
        if smooth:
            ser.smooth = True

        col = s.get("color")
        if col:
            ln = ser.format.line
            ln.color.rgb = to_rgb(col)
            try:
                from pptx.util import Pt as _Pt
                ln.width = _Pt(2)
            except Exception:
                pass

            if markers:
                mk = ser.marker
                try:
                    mk.format.fill.solid()
                    mk.format.fill.fore_color.rgb = to_rgb(col)
                except Exception:
                    pass
                try:
                    mk.format.line.color.rgb = to_rgb(col)
                except Exception:
                    pass
# =========================
# Gantt (chevrons) with left gutter
# =========================
def render_chart_gantt(slide, element):
    # ----- Outer frame -----
    pos = element.get("position", {"x": 1.0, "y": 1.2, "w": 11.3, "h": 5.0})
    x, y, w, h = inches_rect(pos)

    gutter = float(element.get("gutterInches", 1.0))  # label gutter INSIDE the chart
    content_x = x + Inches(gutter)
    content_w = w - Inches(gutter)
    if content_w < Inches(1.0):
        content_w = Inches(1.0)

    # ----- Lanes & lane labels -----
    lanes_raw = element.get("units", {}).get("yRange", {}).get("lanes", []) or []
    lanes = [_shorten_lane(v) for v in lanes_raw] if element.get("shortenLanes", True) else lanes_raw

    lane_count = max(len(lanes), 1)
    lane_height = h / lane_count

    # ----- Time range for items (still proper dates) -----
    units = element.get("units", {}) or {}
    x_range = units.get("xRange", {}) or {}
    t0 = parse_time(x_range.get("t0"))
    t1 = parse_time(x_range.get("t1"))
    total_days = max((t1 - t0).days, 1)

    def x_pos(dt_str: str):
        """Map a date-like string to an x position within content area."""
        dt = parse_time(dt_str)
        d = (dt - t0).days
        return content_x + content_w * (d / total_days)

    def y_pos(lane_index: int):
        return y + lane_height * lane_index

    # lane backgrounds + labels (in gutter)
    for i in range(lane_count):
        yy = y_pos(i)
        # zebra striping
        if i % 2 == 0:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, content_x, yy, content_w, lane_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = to_rgb("#F3F4F6")
            bg.line.fill.background()

        # gutter label
        tb = slide.shapes.add_textbox(
            x,
            yy,
            Inches(gutter) - Inches(0.1),
            lane_height,
        )
        tf = tb.text_frame
        tf.text = lanes[i] if i < len(lanes) else f"Lane {i+1}"
        p = tf.paragraphs[0]
        apply_text_style(
            p,
            {"font": {"size": 12, "bold": True}, "align": "right"},
        )

    # ----- Time grid + top labels (text-only, no date parsing) -----
    grid = element.get("grid") or {}
    quarter_labels = grid.get("quarters") or []    # any text: ["Q1", "Q2"] etc.
    show_labels = bool(grid.get("showLabels", True))
    label_offset = float(grid.get("labelOffsetInches", 0.30))  # inches ABOVE chart
    label_font_size = int(grid.get("labelFontSize", 12))
    draw_axis_line = bool(grid.get("topAxisLine", True))

    n_q = len(quarter_labels)

    if n_q >= 2:
        # vertical guides at evenly spaced boundaries
        for i in range(n_q):
            frac = i / (n_q - 1) if n_q > 1 else 0.0   # 0.0 → 1.0
            xx = content_x + content_w * frac
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, xx, y, Inches(0.018), h
            )
            line.fill.solid()
            line.fill.fore_color.rgb = to_rgb("#E0E0E0")
            line.line.fill.background()

        # top axis line (optional)
        if draw_axis_line:
            top_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                content_x,
                y - Inches(0.06),
                content_w,
                Inches(0.012),
            )
            top_line.fill.solid()
            top_line.fill.fore_color.rgb = to_rgb("#D1D5DB")
            top_line.line.fill.background()

        # labels centered between boundaries, using the provided text as-is
        if show_labels:
            for i in range(n_q - 1):
                frac_mid = (i + 0.5) / (n_q - 1)
                xc = content_x + content_w * frac_mid

                label_text = str(quarter_labels[i])

                tb = slide.shapes.add_textbox(
                    xc - Inches(0.5),
                    y - Inches(label_offset),
                    Inches(1.0),
                    Inches(0.28),
                )
                tf = tb.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = label_text
                apply_text_style(
                    p,
                    {
                        "font": {"size": label_font_size, "bold": False},
                        "align": "center",
                    },
                )

    elif draw_axis_line:
        # no quarters, but still want a top axis line
        top_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            content_x,
            y - Inches(0.06),
            content_w,
            Inches(0.012),
        )
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = to_rgb("#D1D5DB")
        top_line.line.fill.background()

    # ----- Task chevrons -----
    chevron_head = float(element.get("chevronHead", 0.28))

    def draw_chevron(left, top, width, height, fill_hex, text, text_color=DEFAULT_TEXT_COLOR):
        chev = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, left, top, width, height
        )
        try:
            chev.adjustments[0] = chevron_head
        except Exception:
            pass
        chev.fill.solid()
        chev.fill.fore_color.rgb = to_rgb(fill_hex)
        chev.line.fill.background()

        tf = chev.text_frame
        tf.text = text
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.name = DEFAULT_FONT_NAME
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = to_rgb(text_color)
        tf.margin_left = Inches(0.08)
        tf.margin_top = Inches(0.02)
        return chev

    # items
    for it in element.get("items", []):
        start_x = it.get("start", {}).get("x")
        end_x = it.get("end", {}).get("x")

        # if missing start/end, just skip silently (manual input needed)
        if not start_x or not end_x:
            continue

        left = x_pos(start_x)
        right = x_pos(end_x)
        if right < left:
            left, right = right, left
        width = right - left

        lane_idx = int(it.get("start", {}).get("y", 0))
        height = lane_height * float(it.get("size", {}).get("height", 0.6))
        top = y_pos(lane_idx) + (lane_height - height) / 2.0

        style = it.get("style", {}) or {}
        fill = style.get("fill", "#90CAF9")
        text_color = style.get("text", DEFAULT_TEXT_COLOR)
        label = it.get("label", it.get("id", "")) or ""

        draw_chevron(left, top, width, height, fill, label, text_color)

# =========================
# Dispatcher
# =========================
def render_element(slide, element):
    etype = element.get("type")
    if etype == "text":
        render_textbox(slide, element)
    elif etype == "table":
        render_table(slide, element)
    elif etype == "chart":
        subtype = element.get("subtype")
        if subtype == "gantt":
            render_chart_gantt(slide, element)
        elif subtype == "pie":
            render_chart_pie(slide, element)
        elif subtype == "bar":
            render_chart_bar(slide, element)
        elif subtype == "line":
            render_chart_line(slide, element)

# =========================
# Slide helpers (title box + cover slide)
# =========================
def add_title_if_any(slide, title_text: str | None):
    if not title_text:
        return
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(0.35) - Inches(0.25), Inches(11.3), Inches(0.6))
    tf = tb.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    apply_text_style(p, {"font": {"size": 24, "bold": True}, "align": "left"})

def render_cover_slide(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # top band
    band_h = Inches(1.2)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, band_h)
    band.fill.solid()
    band.fill.fore_color.rgb = to_rgb(s.get("accentColor", "#111827"))  # slate-900 as default
    band.line.fill.background()

    # Title (from first heading element)
    title = next((el for el in s.get("elements", []) if el.get("type")=="text" and el.get("variant")=="heading"), None)
    subtitle = next((el for el in s.get("elements", []) if el.get("type")=="text" and el.get("variant") in ("paragraph","rich")), None)

    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), prs.slide_width - Inches(2.0), Inches(1.6))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (title or {}).get("text","")
    p.font.name = DEFAULT_FONT_NAME
    p.font.size = Pt(40); p.font.bold = True
    p.font.color.rgb = to_rgb("#111827")
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), prs.slide_width - Inches(2.0), Inches(0.9))
        stf = sb.text_frame; stf.clear(); stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = subtitle.get("text","")
        sp.font.name = DEFAULT_FONT_NAME
        sp.font.size = Pt(18)
        sp.font.color.rgb = to_rgb("#374151")
        sp.alignment = PP_ALIGN.LEFT

    meta = s.get("meta", {})
    if meta:
        fb = slide.shapes.add_textbox(Inches(1.0), prs.slide_height - Inches(0.9), prs.slide_width - Inches(2.0), Inches(0.5))
        ft = fb.text_frame; ft.clear()
        fp = ft.paragraphs[0]
        fp.text = f'{meta.get("owner","")}  •  {meta.get("date","")}'
        fp.font.size = Pt(12); fp.font.color.rgb = to_rgb("#6B7280"); fp.alignment = PP_ALIGN.LEFT

    return slide

def group_elements_by_slide_flat(elements):
    """Fallback model: flat array with element['slide'] (1-based). Missing ⇒ slide 1."""
    slides = {}
    for el in elements or []:
        idx = int(el.get("slide", 1))
        slides.setdefault(idx, []).append(el)
    return [ {"title": None, "elements": slides[i]} for i in sorted(slides) ]

# =========================
# Deck Builder (Multi-slide)
# =========================
def build_ppt_from_spec(spec: dict, output_filename: str = "spec_demo_output.pptx") -> str:
    prs = Presentation()

    # -------- Slide size (global) --------
    size = (spec.get("slide") or {}).get("size", "16:9")
    if WIDESCREEN_16x9 or size == "16:9":
        # 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        # classic 4:3, if you ever want it
        prs.slide_width = Inches(10.0)
        prs.slide_height = Inches(7.5)

    # -------- Decide slide model --------
    slides_spec = spec.get("slides")
    if not slides_spec:
        # Fallback: flat array with element['slide']
        slides_spec = group_elements_by_slide_flat(spec.get("elements", []))
        if slides_spec and (spec.get("metadata") or {}).get("title"):
            slides_spec[0]["title"] = (spec.get("metadata") or {}).get("title")

    # -------- Build slides --------
    for s in slides_spec:
        if s.get("layout") == "cover":
            slide = render_cover_slide(prs, s)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Only add the big title textbox if there is NO heading element.
            has_heading = any(
                el.get("type") == "text" and el.get("variant") == "heading"
                for el in (s.get("elements") or [])
            )
            if not has_heading:
                add_title_if_any(slide, s.get("title"))

            for el in s.get("elements", []):
                render_element(slide, el)

    out_path = os.path.join(script_dir(), output_filename)
    prs.save(out_path)
    print(f"✅ Saved PowerPoint: {os.path.abspath(out_path)}")
    print(f"   Slides: {len(slides_spec)}")
    return out_path


# --- NEW: compute padded frame if legend would crowd the plot ---
def _padded_chart_frame(position, legend_pos="right", pad_in=0.0):
    """Return (x,y,w,h) Inches with extra space reserved for the legend."""
    from pptx.util import Inches as _In
    x, y, w, h = inches_rect(position)
    pad = _In(max(pad_in, 0.0))
    if legend_pos == "right":
        w = max(_In(1.0), w - pad)          # shrink width
    elif legend_pos == "left":
        x = x + pad; w = max(_In(1.0), w - pad)
    elif legend_pos == "top":
        y = y + pad; h = max(_In(1.0), h - pad)
    elif legend_pos == "bottom":
        h = max(_In(1.0), h - pad)
    return x, y, w, h

def _shorten_lane(s: str) -> str:
    # "Phase 1: Something Long" -> "P1 – Something Long"
    s = s.strip()
    if s.lower().startswith("phase "):
        # pull number if present
        parts = s.split(":", 1)
        head = parts[0]              # e.g., "Phase 1"
        tail = parts[1].strip() if len(parts) > 1 else ""
        num = head.split()[-1]
        short = f"P{num} – {tail}" if tail else f"P{num}"
    else:
        short = s
    # keep it tidy
    return short[:32] + ("…" if len(short) > 32 else "")
# =========================
# SAMPLE: Ten slides (realistic text + safe legends + gantt gutter)
# =========================
def _quarter_label(dt):
    q = ((dt.month - 1) // 3) + 1
    return f"Q{q} {dt.year}"

SAMPLE_TEN_SLIDES = {
      "version": "1.0",
      "slides": [
        {
          "layout": "cover",
          "title": "BankConnect – Internal Employee Management System",
          "accentColor": "#007bff",
          "meta": {
            "owner": "need manual input",
            "date": "need manual input"
          },
          "elements": [
            {
              "type": "text",
              "variant": "heading",
              "text": "BankConnect – Internal Employee Management System",
              "position": { "x": 0.5, "y": 2, "w": 12, "h": 1.5 },
              "style": {
                "font": { "name": "Arial", "size": 44, "bold": True },
                "align": "center",
                "text": "#007bff"
              }
            },
            {
              "type": "text",
              "variant": "paragraph",
              "text": "A unified platform for employee management, communication, and performance tracking",
              "position": { "x": 1, "y": 4, "w": 11, "h": 1 },
              "style": {
                "font": { "name": "Arial", "size": 24 },
                "align": "center"
              }
            }
          ]
        },
        {
      
          "title": "Executive Summary",
          "elements": [
            {
              "type": "text",
              "variant": "heading",
              "text": "Executive Summary",
              "position": { "x": 0.5, "y": 0.5, "w": 12, "h": 1 },
              "style": {
                "font": { "name": "Arial", "size": 36, "bold": True },
                "align": "left",
                "text": "#007bff"
              }
            },
            {
              "type": "text",
              "variant": "bullets",
              "items": [
                { "text": "Objective: Unify employee management across all bank branches", "level": 0 },
                { "text": "Duration: January 2025 – June 2026", "level": 0 },
                { "text": "Users: 500+ internal employees nationwide", "level": 0 },
                { "text": "Target Impact: +30% operational efficiency improvement", "level": 0 }
              ],
              "position": { "x": 0.5, "y": 2, "w": 12, "h": 4 },
              "style": {
                "font": { "name": "Arial", "size": 24 },
                "align": "left"
              }
            }
          ]
        },
        {
      
          "title": "Key Features",
          "elements": [
            {
              "type": "text",
              "variant": "heading",
              "text": "Key Features",
              "position": { "x": 0.5, "y": 0.5, "w": 12, "h": 1 },
              "style": {
                "font": { "name": "Arial", "size": 36, "bold": True },
                "align": "left",
                "text": "#007bff"
              }
            },
            {
              "type": "text",
              "variant": "bullets",
              "items": [
                { "text": "Onboarding: Streamlined onboarding workflow", "level": 0 },
                { "text": "Attendance: Biometric attendance integration", "level": 0 },
                { "text": "Evaluation: 360° feedback system", "level": 0 },
                { "text": "Communication: Secure internal chat and notifications", "level": 0 },
                { "text": "Compliance: Centralized training and certification", "level": 0 },
                { "text": "Analytics: Custom KPI dashboards", "level": 0 }
              ],
              "position": { "x": 0.5, "y": 2, "w": 12, "h": 5 },
              "style": {
                "font": { "name": "Arial", "size": 24 },
                "align": "left"
              }
            }
          ]
        },
        {
      
          "title": "Project Overview",
          "elements": [
            {
              "type": "text",
              "variant": "heading",
              "text": "Project Overview",
              "position": { "x": 0.5, "y": 0.5, "w": 12, "h": 1 },
              "style": {
                "font": { "name": "Arial", "size": 36, "bold": True },
                "align": "left",
                "text": "#007bff"
              }
            },
            {
              "type": "chart",
              "subtype": "bar",
              "title": "Operational Efficiency Improvement",
              "legend": "right",
              "legendPadInches": 1.0,
              "position": { "x": 0.5, "y": 2, "w": 12, "h": 5 },
              "x": { "categories": ["Current", "Target"] },
              "series": [
                { "name": "Efficiency", "data": [70, 100], "color": "#4CAF50" }
              ],
              "options": { "stacked": False, "orientation": "vertical" }
            }
          ]
        },
        {
      
          "title": "Product Roadmap (2025–2026)",
          "elements": [
            {
              "type": "text",
              "variant": "heading",
              "text": "Product Roadmap (2025–2026)",
              "position": { "x": 0.5, "y": 0.5, "w": 12, "h": 1 },
              "style": {
                "font": { "name": "Arial", "size": 36, "bold": True },
                "align": "left",
                "text": "#007bff"
              }
            },
            {
              "type": "chart",
              "subtype": "gantt",
              "title": "Project Timeline",
              "gutterInches": 1.15,
              "shortenLanes": True,
              "position": { "x": 0.5, "y": 2, "w": 12, "h": 5 },
              "units": {
                "xRange": { "t0": "2025-01-01", "t1": "2026-09-30" },
                "yRange": {
                  "lanes": ["Planning", "Design", "Development", "Testing", "Launch"]
                }
              },
              "items": [
                { "id": "1", "label": "Planning", "start": { "x": "2025-01-01", "y": 0 }, "end": { "x": "2025-03-31" }, "style": { "fill": "#4CAF50" } },
                { "id": "2", "label": "Design", "start": { "x": "2025-04-01", "y": 1 }, "end": { "x": "2025-06-30" }, "style": { "fill": "#2196F3" } },
                { "id": "3", "label": "Development", "start": { "x": "2025-07-01", "y": 2 }, "end": { "x": "2026-03-31" }, "style": { "fill": "#FFC107" } },
                { "id": "4", "label": "Testing", "start": { "x": "2026-04-01", "y": 3 }, "end": { "x": "2026-06-30" }, "style": { "fill": "#9C27B0" } },
                { "id": "5", "label": "Launch", "start": { "x": "2026-07-01", "y": 4 }, "end": { "x": "2026-09-30" }, "style": { "fill": "#F44336" } }
              ],
              "grid": {
                "quarters": [
                  "2025 Q1",
                  "2025 Q2",
                  "2025 Q3",
                  "2025 Q4",
                  "2026 Q1",
                  "2026 Q2",
                  "2026 Q3",
                  "2026 Q4"
                ],
                "showLabels": True,
                "labelOffsetInches": 0.2,
                "labelFontSize": 10,
                "topAxisLine": True
              }
            }
          ]
        },
        {
      
          "title": "Risk Management",
          "elements": [
            {
              "type": "text",
              "variant": "heading",
              "text": "Risk Management",
              "position": { "x": 0.5, "y": 0.5, "w": 12, "h": 1 },
              "style": {
                "font": { "name": "Arial", "size": 36, "bold": True },
                "align": "left",
                "text": "#007bff"
              }
            },
            {
              "type": "table",
              "headers": ["Risk", "Impact", "Mitigation"],
              "rows": [
                ["Data Breach", "High", "End-to-end encryption and RBAC policies"],
                ["Delay in Rollout", "Medium", "Agile milestones with regular reviews"],
                ["Adoption Challenges", "Medium", "Training and pilot rollout"],
                ["Integration Errors", "Low", "Dedicated API sandbox testing"]
              ],
              "position": { "x": 0.5, "y": 2, "w": 12, "h": 5 },
              "style": { "align": "left" }
            }
          ]
        },
        {
      
          "title": "Expected Outcomes",
          "elements": [
            {
              "type": "text",
              "variant": "heading",
              "text": "Expected Outcomes",
              "position": { "x": 0.5, "y": 0.5, "w": 12, "h": 1 },
              "style": {
                "font": { "name": "Arial", "size": 36, "bold": True },
                "align": "left",
                "text": "#007bff"
              }
            },
            {
              "type": "text",
              "variant": "bullets",
              "items": [
                { "text": "70% Faster HR Processes", "level": 0 },
                { "text": "80% Reduction in Data Errors", "level": 0 },
                { "text": "20% Increase in Employee Satisfaction", "level": 0 }
              ],
              "position": { "x": 0.5, "y": 2, "w": 12, "h": 4 },
              "style": {
                "font": { "name": "Arial", "size": 24 },
                "align": "left"
              }
            }
          ]
        },
        {
      
          "title": "Closing",
          "elements": [
            {
              "type": "text",
              "variant": "heading",
              "text": "Transforming Banking Operations Through Digital Empowerment",
              "position": { "x": 0.5, "y": 2, "w": 12, "h": 2 },
              "style": {
                "font": { "name": "Arial", "size": 36, "bold": True },
                "align": "center",
                "text": "#007bff"
              }
            },
            {
              "type": "text",
              "variant": "paragraph",
              "text": "📧 contact@bankconnect.com   🌐 www.bankconnect.io",
              "position": { "x": 0.5, "y": 5, "w": 12, "h": 1 },
              "style": {
                "font": { "name": "Arial", "size": 24 },
                "align": "center"
              }
            }
          ]
        }
      ]
}





if __name__ == "__main__":
    # In real use, raw_spec_json comes from your API/body["raw_text"]
    raw_spec_json = json.dumps(SAMPLE_TEN_SLIDES)

    spec = cleaning_JSON(raw_spec_json)
    build_ppt_from_spec(spec, output_filename="roadmap_10slides.pptx")


